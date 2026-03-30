#!/usr/bin/env python3
"""
Excel → PowerPoint slide generator.

For each tab in the Excel file, one slide is produced from the template.
- The text  {{name}}  anywhere in the slide is replaced with the tab name.
- Rows from each tab are appended to the table whose column headers match
  the Excel column headers exactly (non-matching columns are skipped).

Usage:
    python generate_slides.py [excel_file] [template_file] [output_file]

Defaults:
    excel_file    = data.xlsx
    template_file = template.pptx
    output_file   = output.pptx
"""

import copy
import sys
from pathlib import Path

import openpyxl
from lxml import etree
from pptx import Presentation


# ── configuration ─────────────────────────────────────────────────────────────
EXCEL_FILE    = "data.xlsx"
TEMPLATE_FILE = "template.pptx"
OUTPUT_FILE   = "output.pptx"

# Text in the slide that gets replaced with the Excel tab name.
NAME_PLACEHOLDER = "{{name}}"
# Default font size (pt) for inserted data rows when none can be detected from the template.
DEFAULT_FONT_SIZE_PT = 10
# If True, rows that would overflow the slide are placed on new continuation slides.
OVERFLOW_SLIDES = True
# Sheet names to skip entirely (exact match).
EXCLUDE_SHEETS  = ["Bob"]
# Column names where consecutive equal values are merged vertically.
MERGE_COLUMNS   = ["Goal", "Goal2", "Metric"]
# Column names to sort rows by before inserting (ascending, left-to-right priority).
# Rows with None in a sort column are placed last. Excel file is never modified.
SORT_COLUMNS    = ["Goal", "Goal2", "Metric"]
# If True, strip leading/trailing whitespace from cell values before inserting.
STRIP_WHITESPACE = True
# Number of single-line-heights to reserve as blank space at the bottom of each
# slide before triggering overflow.  0 = table may extend to the very bottom edge.
BOTTOM_PADDING_ROWS = 0
# Multiplier applied to each row's estimated height before the overflow comparison.
# Increase above 1.0 to trigger overflow sooner (fewer rows per slide / more conservative).
# Decrease below 1.0 to trigger overflow later (more rows per slide / more permissive).
OVERFLOW_SENSITIVITY = 0.75
# ──────────────────────────────────────────────────────────────────────────────


NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Cell margin matching the template's tcPr (marL=marR=45720, marT=marB default=45720).
CELL_MARGIN_EMU = 45720   # 0.05 in = 45 720 EMU per side

# Line-height multiplier used in row-height estimation.
# Calibrated from actual PowerPoint-exported PNGs (1920×1080, scale 6350 EMU/px):
#   single-line Arial 11 pt data row ≈ 56 px = 355 600 EMU
#   → text portion = 355 600 − 2×45 720 = 264 160 EMU
#   → 264 160 / (11 × 12 700) = 1.890 → round up to 1.95 for safety margin.
LINE_HEIGHT_MULTIPLIER = 1.95

# Per-character width weights (em units) for proportional-font wrap estimation.
# Calibrated for Arial; applies well to similar sans-serif fonts.
_NARROW_CHARS = frozenset(' ifjlrtI!|:;.,\'"`1()')
_WIDE_CHARS   = frozenset('MWmw')


# ── helpers ───────────────────────────────────────────────────────────────────

def detect_font_size(tr):
    """Return font size in points from the first <a:rPr sz=...> found in *tr*, or None."""
    for rPr in tr.iter(f"{{{NS}}}rPr"):
        sz = rPr.get("sz")
        if sz:
            return int(sz) // 100
    return None


def detect_para_spacing(trs):
    """
    Return total paragraph spacing in EMU from template data rows *trs*.

    Sums the maximum spcBef and spcAft (in EMU) found across all <a:pPr>
    elements.  Only fixed-point spacing (<a:spcPts>) is handled; percentage-
    based spacing contributes 0.  Returns 0 when no explicit spacing is set.

    Used to correct multi-line height estimates: PowerPoint applies spcBef/
    spcAft once per paragraph, not once per visual wrapped line, so charging
    the full LINE_HEIGHT_EMU per additional wrapped line over-estimates height.
    """
    spc_bef = spc_aft = 0
    for tr in trs:
        for pPr in tr.iter(f"{{{NS}}}pPr"):
            for tag, is_bef in (
                (f"{{{NS}}}spcBef", True),
                (f"{{{NS}}}spcAft", False),
            ):
                parent = pPr.find(tag)
                if parent is None:
                    continue
                pts_el = parent.find(f"{{{NS}}}spcPts")
                if pts_el is not None:
                    # spcPts val is in 1/100th of a point; 1/100 pt = 127 EMU.
                    emu = int(pts_el.get("val", 0)) * 127
                    if is_bef:
                        spc_bef = max(spc_bef, emu)
                    else:
                        spc_aft = max(spc_aft, emu)
    return spc_bef + spc_aft


def _word_wrap_lines(text, usable_w_emu, font_size_pt):
    """
    Estimate the number of visual lines for *text* rendered in Arial
    *font_size_pt* within *usable_w_emu* EMU of usable horizontal space.

    Improvements over the previous naive ceil(len/cpl) approach:
    - Breaks only at word boundaries (spaces), matching PowerPoint behaviour.
    - Uses per-character width weights (_NARROW_CHARS / _WIDE_CHARS) instead
      of a uniform 0.5-em average for more accurate CPL on business text.
    - Words wider than the column are broken at character boundaries.
    """
    if not text:
        return 1
    words = text.split()
    if not words:
        return 1

    em      = font_size_pt * 12700   # 1 em in EMU at this font size
    space_w = 0.34 * em

    def _word_w(w):
        total = 0.0
        for ch in w:
            if ch in _NARROW_CHARS:
                total += 0.34
            elif ch in _WIDE_CHARS:
                total += 0.72
            elif ch.isupper():
                total += 0.65
            else:
                total += 0.52
        return total * em

    lines = 1
    cur   = 0.0

    for word in words:
        ww = _word_w(word)
        if cur == 0:
            if ww > usable_w_emu:
                # Single word wider than the column — force character breaks.
                q, r = divmod(ww, usable_w_emu)
                lines += int(q)
                cur    = r if r > 0 else usable_w_emu
            else:
                cur = ww
        elif cur + space_w + ww <= usable_w_emu:
            cur += space_w + ww
        else:
            lines += 1
            if ww > usable_w_emu:
                q, r = divmod(ww, usable_w_emu)
                lines += int(q)
                cur    = r if r > 0 else usable_w_emu
            else:
                cur = ww

    return lines


def get_table_shape(slide):
    """Return the first table shape found on *slide*, or None."""
    for shape in slide.shapes:
        if shape.has_table:
            return shape
    return None


def cell_text(cell):
    """Return stripped text content of a table cell."""
    return cell.text_frame.text.strip()


def estimate_row_height(row_data, col_map, col_widths, font_size, para_spacing_emu=0):
    """
    Estimate the rendered height (EMU) of a table row from its text content.

    Uses word-boundary wrapping and per-character width estimates (via
    _word_wrap_lines) to count visual lines per cell, then returns the height
    for the tallest cell plus standard top/bottom cell margins.

    para_spacing_emu : spcBef + spcAft detected from template (EMU).
        PowerPoint applies paragraph spacing once per paragraph, not once per
        visual wrapped line.  For multi-line cells, each line beyond the first
        should NOT be charged the full LINE_HEIGHT_EMU (which includes spacing);
        only the raw line height (LINE_HEIGHT_EMU − para_spacing_emu) is added.

    col_widths : {pptx_col_index: column_width_in_EMU}
    """
    SIDE_MARGIN_EMU = CELL_MARGIN_EMU * 2       # left + right
    VERT_MARGIN_EMU = CELL_MARGIN_EMU * 2       # top  + bottom
    LINE_HEIGHT_EMU = int(font_size * LINE_HEIGHT_MULTIPLIER * 12700)
    # Additional wrapped lines add only the raw line height — para spacing
    # (spcBef/spcAft) is not repeated for each visual line in the same paragraph.
    PER_EXTRA_LINE  = max(LINE_HEIGHT_EMU - para_spacing_emu, 1)

    max_lines = 1
    for col_idx, col_name in col_map.items():
        text = str(row_data.get(col_name) or "")
        if not text:
            continue
        usable_w  = max((col_widths.get(col_idx) or 0) - SIDE_MARGIN_EMU,
                        font_size * 12700)       # at least one em wide
        lines     = _word_wrap_lines(text, usable_w, font_size)
        max_lines = max(max_lines, lines)

    # First line: full LINE_HEIGHT_EMU (calibrated value, includes any para spacing once).
    # Each additional line: only the raw per-line height (para spacing not repeated).
    return LINE_HEIGHT_EMU + (max_lines - 1) * PER_EXTRA_LINE + VERT_MARGIN_EMU


def append_data_row(table, col_map, row_data, font_size=DEFAULT_FONT_SIZE_PT):
    """
    Build a brand-new <a:tr> element and append it to *table*.

    Rows are constructed from scratch rather than copied from an existing row.
    Copying any existing row (including the header) inherits cell properties —
    fill colours, run-formatting flags, and sometimes implicit lock attributes —
    that cause PowerPoint to render the pasted cells as non-editable.
    A clean minimal row avoids all of that; the table's built-in style still
    applies banded-row colours automatically based on each row's position.

    col_map  : {pptx_column_index: excel_column_name}
    row_data : {excel_column_name: value}
    """
    tbl       = table._tbl
    header_tr = tbl.findall(f"{{{NS}}}tr")[0]
    num_cols  = len(header_tr.findall(f"{{{NS}}}tc"))

    new_tr = etree.Element(f"{{{NS}}}tr")
    new_tr.set("h", "0")   # 0 = auto-height; PowerPoint sizes each row to its content

    for col_idx in range(num_cols):
        col_name = col_map.get(col_idx)
        if col_name is not None:
            val  = row_data.get(col_name)
            text = "" if val is None else str(val)
        else:
            text = ""

        tc     = etree.SubElement(new_tr, f"{{{NS}}}tc")
        txBody = etree.SubElement(tc, f"{{{NS}}}txBody")
        etree.SubElement(txBody, f"{{{NS}}}bodyPr")
        etree.SubElement(txBody, f"{{{NS}}}lstStyle")
        p      = etree.SubElement(txBody, f"{{{NS}}}p")

        if text:
            r   = etree.SubElement(p, f"{{{NS}}}r")
            rPr = etree.SubElement(r, f"{{{NS}}}rPr")
            rPr.set("lang", "en-US")
            rPr.set("dirty", "0")
            rPr.set("sz", str(font_size * 100))
            etree.SubElement(rPr, f"{{{NS}}}latin").set("typeface", "Arial")
            t   = etree.SubElement(r, f"{{{NS}}}t")
            t.text = text

        # Always write endParaRPr so empty cells use the template font
        # (Arial / detected size) instead of falling back to PowerPoint's
        # table-style default (Aptos 18pt).
        endPr = etree.SubElement(p, f"{{{NS}}}endParaRPr")
        endPr.set("lang", "en-US")
        endPr.set("dirty", "0")
        endPr.set("sz", str(font_size * 100))
        etree.SubElement(endPr, f"{{{NS}}}latin").set("typeface", "Arial")

        tc_pr = etree.SubElement(tc, f"{{{NS}}}tcPr")
        for attr in ("marL", "marR", "marT", "marB"):
            tc_pr.set(attr, str(CELL_MARGIN_EMU))

    tbl.append(new_tr)


def replace_placeholder(slide, placeholder, value):
    """Replace every occurrence of *placeholder* in all text runs on *slide*."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)


def make_slide_from_template(prs, original_sp_tree, layout):
    """
    Add a new slide to *prs* and populate it with a fresh copy of
    *original_sp_tree* (the shape tree snapshotted from the template slide).
    """
    new_slide = prs.slides.add_slide(layout)
    sp_tree = new_slide.shapes._spTree

    # Remove the auto-generated placeholder shapes (keep first 2 fixed children:
    # nvGrpSpPr and grpSpPr).
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)

    # Paste a fresh copy of the original template shapes.
    for child in list(original_sp_tree)[2:]:
        sp_tree.append(copy.deepcopy(child))

    return new_slide


def apply_vertical_merges(table, col_map, merge_cols):
    """
    Vertically merge consecutive cells in the specified columns.

    Non-empty cells: merge runs of equal values as usual.
    Empty cells in a column that has a MERGE_COLUMN to its left: merge
    sub-groups that align with the left merge-column's groupings, so blank
    cells visually track the parent group to their left.
    Empty cells in the leftmost merge column: sub-divide by the immediate
    right column's value groups (fallback behaviour).

    Group IDs are pre-computed from raw text BEFORE any merges fire, because
    merge operations clear continuation-cell text and would corrupt later reads.

    col_map   : {pptx_col_index: excel_col_name}
    merge_cols: list of excel column names to merge
    """
    name_to_idx = {name: idx for idx, name in col_map.items()}
    tbl         = table._tbl
    data_trs    = tbl.findall(f"{{{NS}}}tr")[1:]   # skip header row

    if not data_trs:
        return

    num_tcs = len(data_trs[0].findall(f"{{{NS}}}tc"))

    def _do_merge(cells, start, end):
        """Apply rowSpan on cells[start] and vMerge on cells[start+1:end]."""
        if end - start < 2:
            return
        cells[start][1].set("rowSpan", str(end - start))
        for k in range(start + 1, end):
            cont_tc = cells[k][1]
            cont_tc.set("vMerge", "1")
            txBody = cont_tc.find(f"{{{NS}}}txBody")
            if txBody is not None:
                for p in txBody.findall(f"{{{NS}}}p"):
                    txBody.remove(p)
                etree.SubElement(txBody, f"{{{NS}}}p")

    def _read_col_vals(col_idx):
        vals = []
        for tr in data_trs:
            tcs   = tr.findall(f"{{{NS}}}tc")
            t_el  = tcs[col_idx].find(f".//{{{NS}}}t") if col_idx < len(tcs) else None
            vals.append((t_el.text or "").strip() if t_el is not None else "")
        return vals

    def _group_ids(vals):
        """
        Return a list where each entry is the start-index of that row's
        consecutive-equal-value run.  Two rows share a group iff they have
        the same group ID.
        """
        ids = [0] * len(vals)
        i = 0
        while i < len(vals):
            j = i + 1
            while j < len(vals) and vals[j] == vals[i]:
                j += 1
            for k in range(i, j):
                ids[k] = i
            i = j
        return ids

    # Identify merge-column pptx indices, sorted left → right.
    merge_col_idxs = sorted(
        idx for name in merge_cols
        for idx in [name_to_idx.get(name)] if idx is not None
    )
    merge_col_set = set(merge_col_idxs)

    # Pre-read every merge column's raw values and compute group IDs BEFORE
    # any merges fire (merges clear continuation-cell text, corrupting reads).
    pre_group_ids = {
        col_idx: _group_ids(_read_col_vals(col_idx))
        for col_idx in merge_col_set
    }

    for col_name in merge_cols:
        col_idx = name_to_idx.get(col_name)
        if col_idx is None:
            continue   # column not present in this table — skip silently

        # Collect [text, <a:tc>] for every data row in this column.
        cells = []
        for tr in data_trs:
            tcs   = tr.findall(f"{{{NS}}}tc")
            tc    = tcs[col_idx] if col_idx < len(tcs) else None
            t_el  = tc.find(f".//{{{NS}}}t") if tc is not None else None
            value = (t_el.text or "").strip() if t_el is not None else ""
            cells.append([value, tc])

        # Choose the reference grouping for empty-run sub-division.
        # Prefer the nearest MERGE_COLUMN to the left so that empty cells
        # track their parent group (the column that "owns" this hierarchy level).
        # Fall back to the immediate right column for the leftmost merge column.
        left_merge_idx = max(
            (idx for idx in merge_col_set if idx < col_idx), default=None
        )
        if left_merge_idx is not None:
            ref_gids = pre_group_ids[left_merge_idx]
        else:
            right_col = col_idx + 1
            ref_gids  = _group_ids(_read_col_vals(right_col)) if right_col < num_tcs else []

        i = 0
        while i < len(cells):
            value = cells[i][0]
            j = i + 1
            while j < len(cells) and cells[j][0] == value:
                j += 1

            if value:
                # Non-empty run: merge the whole group if span > 1.
                _do_merge(cells, i, j)
            elif ref_gids and (j - i) > 1:
                # Empty run: sub-divide by the reference column's groups
                # and merge each sub-group independently.
                k = i
                while k < j:
                    gid = ref_gids[k]
                    m   = k + 1
                    while m < j and ref_gids[m] == gid:
                        m += 1
                    _do_merge(cells, k, m)
                    k = m

            i = j


# ── main logic ────────────────────────────────────────────────────────────────

def export_slide_pngs(layouts, slide_w, slide_h, output_path):
    """
    Write one PNG per slide showing the estimated row layout.

    Each PNG draws the table as coloured row-bands with a red dashed line
    at the overflow boundary (fill_height) so overflow can be spotted at a
    glance without opening PowerPoint.

    layouts : list of dicts —
        name, table_left, table_top, table_width,
        header_height, fill_height, rows (list of estimated row heights in EMU)
    """
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        print("  [warn] Pillow not available; skipping PNG export.")
        return

    stem   = str(output_path).rsplit(".", 1)[0]
    SCALE  = 1920 / slide_w          # normalise to 1920 px wide
    W, H   = 1920, int(slide_h * SCALE)

    for idx, layout in enumerate(layouts):
        img  = Image.new("RGB", (W, H), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        tx = int(layout["table_left"]  * SCALE)
        tw = int(layout["table_width"] * SCALE)
        ty = int(layout["table_top"]   * SCALE)

        # Header band
        hy = ty + int(layout["header_height"] * SCALE)
        draw.rectangle([(tx, ty), (tx + tw - 1, hy - 1)], fill=(30, 90, 140))

        # Data row bands (alternating tints)
        y = hy
        for j, rh in enumerate(layout["rows"]):
            rh_px = max(1, int(rh * SCALE))
            fill  = (200, 220, 235) if j % 2 == 0 else (230, 230, 230)
            draw.rectangle([(tx, y), (tx + tw - 1, y + rh_px - 1)], fill=fill,
                           outline=(180, 180, 180), width=1)
            y += rh_px

        # Fill-height boundary (red dashed)
        fy = int(layout["fill_height"] * SCALE)
        for x in range(0, W, 24):
            draw.line([(x, fy), (min(x + 12, W), fy)], fill=(220, 30, 30), width=2)

        # Slide bottom boundary (thin black line)
        draw.line([(0, H - 1), (W - 1, H - 1)], fill=(0, 0, 0), width=3)

        # Flag overflow: any row that extends below fill_height
        if y > fy:
            overflow_y = max(fy, ty)
            draw.rectangle([(tx, overflow_y), (tx + tw - 1, min(y, H) - 1)],
                           fill=(255, 100, 100, 180))
            draw.text((tx + 4, overflow_y + 4), "OVERFLOW", fill=(180, 0, 0))

        png_path = f"{stem}_slide{idx + 1:02d}.png"
        img.save(png_path)
        print(f"  PNG saved: {png_path}")


def process(
    excel_path,
    template_path,
    output_path,
    placeholder=NAME_PLACEHOLDER,
    overflow_slides=OVERFLOW_SLIDES,
    exclude_sheets=None,
    merge_columns=None,
    sort_columns=None,
    strip_whitespace=STRIP_WHITESPACE,
    bottom_padding_rows=BOTTOM_PADDING_ROWS,
    overflow_sensitivity=OVERFLOW_SENSITIVITY,
):
    exclude_sheets = set(exclude_sheets or [])
    merge_columns  = list(merge_columns  or [])
    sort_columns   = list(sort_columns   or [])

    wb  = openpyxl.load_workbook(excel_path)
    prs = Presentation(template_path)

    if not prs.slides:
        sys.exit("Error: template contains no slides.")

    # Snapshot the original template slide *before* we mutate anything.
    tmpl_slide       = prs.slides[0]
    original_sp_tree = copy.deepcopy(tmpl_slide.shapes._spTree)
    tmpl_layout      = tmpl_slide.slide_layout

    sheets = wb.sheetnames
    print(f"Sheets found: {', '.join(sheets)}\n")

    slides_created   = []
    slide_layouts    = []   # collected for PNG export
    first_slide_used = False

    for sheet_name in sheets:
        if sheet_name in exclude_sheets:
            print(f"[SKIP] '{sheet_name}': in exclude list.")
            continue

        ws = wb[sheet_name]

        # Read column headers from row 1 (stop at first blank cell).
        headers = []
        for cell in ws[1]:
            if cell.value is None:
                break
            headers.append(str(cell.value).strip())

        if not headers:
            print(f"[SKIP] '{sheet_name}': no headers found in row 1.")
            continue

        # Read data rows (skip fully empty rows).
        data_rows = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            if all(v is None for v in row):
                continue
            data_rows.append({
                headers[c]: row[c]
                for c in range(min(len(headers), len(row)))
            })

        # Strip leading/trailing whitespace from all string values upfront,
        # so that sort keys and overflow height estimates use clean text.
        if strip_whitespace:
            data_rows = [
                {k: v.strip() if isinstance(v, str) else v for k, v in row.items()}
                for row in data_rows
            ]

        # Sort rows in-memory (Excel file is never modified).
        sort_cols = [c for c in sort_columns if c in headers]
        if sort_cols:
            data_rows.sort(key=lambda row: tuple(
                (0 if row.get(c) is not None else 1, str(row.get(c) or "").lower())
                for c in sort_cols
            ))

        print(f"Sheet '{sheet_name}': {len(data_rows)} data row(s), "
              f"columns: {headers}")

        # Get or create the first slide for this sheet.
        if not first_slide_used:
            slide            = prs.slides[0]
            first_slide_used = True
        else:
            slide = make_slide_from_template(prs, original_sp_tree, tmpl_layout)

        replace_placeholder(slide, placeholder, sheet_name)

        # Locate the table shape (shape needed for .top in overflow calc).
        table_shape = get_table_shape(slide)
        if table_shape is None:
            print(f"  [WARN] No table found on slide — skipping '{sheet_name}'.")
            continue
        table = table_shape.table

        # Map: PowerPoint column index → Excel column name (exact-match only).
        pptx_cols = [
            cell_text(table.rows[0].cells[c])
            for c in range(len(table.rows[0].cells))
        ]
        col_map = {c: name for c, name in enumerate(pptx_cols) if name in headers}

        if not col_map:
            print(f"  [WARN] No column names match between Excel and the table.")
            print(f"         Table headers : {pptx_cols}")
            print(f"         Excel headers : {headers}")
            continue

        matched = [col_map[c] for c in sorted(col_map)]
        skipped = [h for h in headers if h not in matched]
        print(f"  Matched : {matched}")
        if skipped:
            print(f"  Skipped : {skipped}")

        # Remove template placeholder rows; detect font size first.
        tbl          = table._tbl
        existing_trs = tbl.findall(f"{{{NS}}}tr")
        font_size    = DEFAULT_FONT_SIZE_PT
        for tr in existing_trs[1:]:
            detected = detect_font_size(tr)
            if detected:
                font_size = detected
                break

        # Detect paragraph spacing from template rows before removing them.
        # spcBef/spcAft apply once per paragraph, not per visual wrapped line;
        # knowing the value prevents estimate_row_height from double-counting
        # paragraph spacing on each additional wrapped line in a multi-line cell.
        para_spacing_emu = detect_para_spacing(existing_trs[1:])

        # Header row height for overflow accounting (read from template).
        hdr_el     = existing_trs[0]
        hdr_h      = hdr_el.get("h", "")
        fallback_h = int(font_size * LINE_HEIGHT_MULTIPLIER * 12700) + 2 * CELL_MARGIN_EMU
        try:
            header_height = int(hdr_h) if hdr_h and int(hdr_h) > 0 else fallback_h
        except (ValueError, TypeError):
            header_height = fallback_h

        for tr in existing_trs[1:]:
            tbl.remove(tr)

        # Column widths (EMU) used by estimate_row_height for overflow detection.
        col_widths = {col_idx: table.columns[col_idx].width for col_idx in col_map}

        # Bottom padding: reserve N full row-heights above the slide edge.
        # A full row = text line height + top/bottom cell margins, matching
        # what estimate_row_height returns for a single-line row.
        line_height_emu   = int(font_size * LINE_HEIGHT_MULTIPLIER * 12700)
        single_row_emu    = line_height_emu + CELL_MARGIN_EMU * 2
        bottom_pad_emu    = bottom_padding_rows * single_row_emu
        slide_fill_height = prs.slide_height - bottom_pad_emu

        # Pre-compute which (row_index, col_name) pairs will be rendered as
        # merged/spanned cells.  In a rowSpan group, PowerPoint sets each
        # individual row's height from its non-spanning cells; the spanning
        # cell's text flows across the combined height of the group.  So we
        # exclude those cells from per-row height estimation.
        merge_spanned: set = set()
        if overflow_slides and merge_columns:
            for col_name in merge_columns:
                if col_name not in headers:
                    continue
                i = 0
                while i < len(data_rows):
                    val = data_rows[i].get(col_name)
                    j   = i + 1
                    while j < len(data_rows) and data_rows[j].get(col_name) == val:
                        j += 1
                    if (j - i) > 1 and val:
                        for k in range(i, j):
                            merge_spanned.add((k, col_name))
                    i = j

        current_table     = table
        current_table_top = table_shape.top if overflow_slides else 0
        current_h         = current_table_top + header_height

        slides_created.append(sheet_name)
        rows_on_current_slide = 0
        current_layout = {
            "name": sheet_name, "rows": [],
            "table_left": table_shape.left, "table_top": current_table_top,
            "table_width": table_shape.width, "header_height": header_height,
            "fill_height": slide_fill_height if overflow_slides else prs.slide_height,
        }
        slide_layouts.append(current_layout)

        # Insert data rows, spilling onto continuation slides when needed.
        for row_idx, row_data in enumerate(data_rows):
            if overflow_slides:
                if merge_spanned:
                    effective = {k: (None if (row_idx, k) in merge_spanned else v)
                                 for k, v in row_data.items()}
                else:
                    effective = row_data
                row_h = estimate_row_height(effective, col_map, col_widths, font_size, para_spacing_emu) * overflow_sensitivity
                # Guard: only overflow when there is already at least one data row,
                # preventing an infinite loop if a single row exceeds slide height.
                if rows_on_current_slide >= 1 and current_h + row_h > slide_fill_height:
                    if merge_columns:
                        apply_vertical_merges(current_table, col_map, merge_columns)

                    cont_label    = f"{sheet_name} (cont.)"
                    ovf_slide     = make_slide_from_template(prs, original_sp_tree, tmpl_layout)
                    replace_placeholder(ovf_slide, placeholder, cont_label)
                    slides_created.append(cont_label)

                    ovf_shape         = get_table_shape(ovf_slide)
                    current_table     = ovf_shape.table
                    current_table_top = ovf_shape.top
                    ovf_tbl           = current_table._tbl
                    for tr in ovf_tbl.findall(f"{{{NS}}}tr")[1:]:
                        ovf_tbl.remove(tr)

                    current_h             = current_table_top + header_height
                    rows_on_current_slide = 0
                    current_layout = {
                        "name": cont_label, "rows": [],
                        "table_left": ovf_shape.left, "table_top": current_table_top,
                        "table_width": ovf_shape.width, "header_height": header_height,
                        "fill_height": slide_fill_height,
                    }
                    slide_layouts.append(current_layout)

                append_data_row(current_table, col_map, row_data, font_size)
                current_h             += row_h
                rows_on_current_slide += 1
                current_layout["rows"].append(row_h)
            else:
                append_data_row(current_table, col_map, row_data, font_size)

        # Apply merges to the final (or only) table for this sheet.
        if merge_columns:
            apply_vertical_merges(current_table, col_map, merge_columns)

    prs.save(output_path)
    print(f"\nDone — {len(slides_created)} slide(s) created: {', '.join(slides_created)}")
    print(f"Output saved to: {output_path}")
    print("\nExporting slide PNGs...")
    export_slide_pngs(slide_layouts, prs.slide_width, prs.slide_height, output_path)


# ── entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    args = sys.argv[1:]
    excel    = args[0] if len(args) > 0 else EXCEL_FILE
    template = args[1] if len(args) > 1 else TEMPLATE_FILE
    output   = args[2] if len(args) > 2 else OUTPUT_FILE

    for path, label in [(excel, "Excel"), (template, "Template")]:
        if not Path(path).exists():
            sys.exit(f"Error: {label} file not found: '{path}'")

    process(
        excel, template, output,
        overflow_slides=OVERFLOW_SLIDES,
        exclude_sheets=EXCLUDE_SHEETS,
        merge_columns=MERGE_COLUMNS,
        sort_columns=SORT_COLUMNS,
        strip_whitespace=STRIP_WHITESPACE,
        bottom_padding_rows=BOTTOM_PADDING_ROWS,
        overflow_sensitivity=OVERFLOW_SENSITIVITY,
    )
